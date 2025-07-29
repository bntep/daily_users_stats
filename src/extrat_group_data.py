# -*- coding: utf-8 -*-
import os
import pandas as pd
import sys
import argparse
from pathlib import Path
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
sys.path.append(str(Path(os.getcwd())))
from utils.dbclient.DatabaseClient import DbConnector
from module.env import *



parser = argparse.ArgumentParser(description="This script creates a ppt slides with statistics for daily users database. \n \
                                 Param: -l or --labo(or list of laboratories)")

parser.add_argument('--year', '-y',
                    help="a year or --year(or list of years)")

parser.add_argument('--labo', '-l',
                    help="the name of the laboratory", required=False)

args = parser.parse_args()


#supprime les caractères et les espaces inattendus et met en minuscules tous les noms d'institutions
def normalize(institution):
    return str(institution).strip().lower()
    
def extract_data(folder_path, institution_list):

    select_csv_columns = {        
        #'raw_data_stats_daily_subscription.csv': ['id_user', 'labo_name', 'statut', 'date_create_year', 'date_last_access_year'],  # type: ignore
        'raw_data_stats_daily_users_with_subscription.csv': ['id_user', 'user_name', 'labo_name', 'statut', 'date_last_access', 'date_created'] # type: ignore
        }
    

    #creates institution dictionary
    #crée un dictionnaire institutionnel
    institution_data = {}

    for file in os.listdir(folder_path):
        if file.endswith(".csv"):
            try:
                #reads csv files and formats to ensure path works across many operating systems
                #lit les fichiers et formats CSV pour garantir que le chemin fonctionne sur de nombreux systèmes d'exploitation
                df = pd.read_csv(os.path.join(folder_path, file), delimiter = "|")
                
                
                # Normalizes institution names to handle case insensitivity and unnecessary characters
                normalized_institution_list = []
                normalized_institution_list = [normalize(inst) for inst in institution_list]
                
                #creates new institution list that ignores capitalization and unnecessary characters
                #crée une nouvelle liste d'institutions qui ignore les majuscules et les caractères inutiles
                institution_list = normalized_institution_list

                #finds if any part of each row contains an institution
                #détermine si une partie de chaque ligne contient une institution
                match_institution = df.apply(lambda row: any(any(normalized_name in normalize(cell)for normalized_name in normalized_institution_list)for cell in row),axis=1)

                #filters dataframe to include only rows where match_institution is True
                #filtre le dataframe pour inclure uniquement les lignes où match_institution est True
                if not match_institution.any():
                    continue
                filtered_df = df[match_institution]

                # If custom columns are defined for this file, use only those columns (if they exist in the DataFrame)
                if file in select_csv_columns:
                    columns_to_keep = [col for col in select_csv_columns[file] if col in filtered_df.columns]
                    filtered_df = filtered_df[columns_to_keep]
                    filtered_df.drop_duplicates(inplace=True)
                    filtered_df.reset_index(drop=True, inplace=True)

                for institution in institution_list:
                    #filters only the rows that contain the current institution
                    #filtre uniquement les lignes qui contiennent l'institution actuelle
                    institution_rows = filtered_df[filtered_df.apply(lambda row: any(institution in normalize(str(cell)) for cell in row), axis=1)].copy()

                    # Attempt to convert year-like float values to integers (e.g., 2022.0 -> 2022)
                    # Essaie de convertir les valeurs de type float en années en entiers (par exemple, 2022.0 -> 2022)
                    for col in institution_rows.columns:
                        if institution_rows[col].dtype == 'float64':
                            # Only convert if all values are integers in float form
                            if institution_rows[col].dropna().apply(float.is_integer).all():
                                institution_rows[col] = institution_rows[col].astype('Int64')

                    #skip institution if not in file
                    #ignorer l'établissement s'il n'est pas dans le fichier
                    if not institution_rows.empty:
                        
                        #creates new list to store this institutions data if it does not alredy exist in institution_data dictionary
                        #crée une nouvelle liste pour stocker les données de cette institution si elles n'existent pas déjà dans le dictionnaire de données de l'institution
                        if institution not in institution_data:
                            institution_data[institution] = []
        
                        institution_data[institution].append({
                            "file": file,
                            "columns": institution_rows.columns.tolist(),
                            "data": institution_rows
                        })

                        # drop duplicates and reset index for the last entry in the institution_data list
                        # This ensures that each institution's data is unique and properly indexed  
                        institution_data[institution][-1]["data"].drop_duplicates(inplace=True)
                        institution_data[institution][-1]["data"].reset_index(drop=True, inplace=True)

            except Exception as e:
                # Handle exceptions for individual files, but continue processing others        
                print(f"Error processing {file}: {e}")
    return institution_data

def chunk_data_rows(rows, chunk_size):
    """Yield successive chunks from list of rows."""
    for i in range(0, len(rows), chunk_size):
        yield rows[i:i + chunk_size]

def create_institutional_pptx(institution_data, presentation, max_rows_per_slide=10):
    replace_existing = True  # or False based on user input

    if replace_existing or not os.path.exists(presentation):
        prs = Presentation()  # start fresh
    else:
        prs = Presentation(presentation)  # load and append

    for institution, blocks in institution_data.items():
        for data in blocks:
            full_rows = data["data"].to_dict(orient="records")
            row_chunks = list(chunk_data_rows(full_rows, max_rows_per_slide))
            
            for idx, chunk in enumerate(row_chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                # define the text title policy
                
                slide.shapes.title.text = f"{institution} ({idx + 1}/{len(row_chunks)})"
                title_font = title_shape.text_frame.paragraphs[0].font
                title_font.name = 'Calibri'
                title_font.size = Pt(24)
                title_font.bold = True
                textbox = slide.placeholders[1]               
                tf = textbox.text_frame               
                tf.clear()

                # File name header
                p = tf.paragraphs[0]
                p.text = f"{data['file']}"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.name = 'Bodoni MT Condensed'

                # Column headers
                header_p = tf.add_paragraph()
                header_p.text = " | ".join(data["columns"])
                header_p.font.bold = True
                header_p.font.size = Pt(12)
                header_p.font.name = 'Bodoni MT Condensed'
                header_p.level = 1

                # Add each row as values only, with statut if file is 'raw_data_stats_daily_users.csv'
                for row_data in chunk:
                    values = " | ".join([str(row_data[col]) for col in data["columns"]])
                    p = tf.add_paragraph()
                    p.text = values
                    p.font.size = Pt(12)
                    p.font.name = 'Bodoni MT Condensed'
                    p.level = 1


    prs.save(presentation)
    return presentation


def main():
    req_extract_stats_labo_nom = "select distinct nom_groupe_labo from statistique_requete;"
    
    df_nom_groupe_labo = DbConnector('yakari', echo = True).execute_query(req_extract_stats_labo_nom)
    
    institution_list = df_nom_groupe_labo.values.flatten().tolist()

    if args.labo:
        # If a specific laboratory is provided, filter the institution list
        institution_list = [normalize(args.labo)]

    #data_folder_path = Path("/mnt/c/Users/akash/Documents/DEV_python/daily_users_stats/user_data_csv")
    data_folder_path = Path("/home/groups/daily/travail/Akash/DEV_python/daily_users_stats/user_data_csv")
    institution_data =  extract_data(data_folder_path, institution_list)

    if args.labo:
        # If a specific laboratory is provided, filter the institution data
        institution_data = {k: v for k, v in institution_data.items() if normalize(k) == normalize(args.labo)}
        output_pptx = Path(f"/home/groups/daily/travail/Akash/DEV_python/daily_users_stats/akash_output/{args.labo}.pptx")
        create_institutional_pptx(institution_data, presentation=output_pptx)
    
    else:
        output_pptx = Path("/home/groups/daily/travail/Akash/DEV_python/daily_users_stats/akash_output/institution_summary.pptx")
        #output_pptx = Path("/mnt/c/Users/akash/Documents/DEV_python/daily_users_stats/akash_output/institution_summary.pptx")
        create_institutional_pptx(institution_data, presentation=output_pptx)



if __name__ == "__main__":
    main()
  